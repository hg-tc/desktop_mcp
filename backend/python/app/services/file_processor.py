"""
文件处理服务 - 简化版
支持多种文件格式的解析和内容提取
"""
import os
import logging
import json
from typing import Dict, List, Any
from pathlib import Path

logger = logging.getLogger(__name__)

# 尝试导入文件处理库（可选依赖）
try:
    import fitz  # PyMuPDF
    HAS_PYMUPDF = True
except ImportError:
    HAS_PYMUPDF = False
    logger.warning("PyMuPDF 未安装，PDF 处理功能不可用")

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    logger.warning("python-docx 未安装，Word 处理功能不可用")

try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False
    logger.warning("openpyxl 未安装，Excel 处理功能不可用")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    logger.warning("python-pptx 未安装，PPT 处理功能不可用")

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    logger.warning("Pillow 未安装，图片处理功能受限")


class FileProcessor:
    """文件处理器"""

    def __init__(self, upload_dir: str = "uploads"):
        self.upload_dir = Path(upload_dir)
        self.upload_dir.mkdir(parents=True, exist_ok=True)

    def process_file(self, file_path: str) -> Dict[str, Any]:
        """
        处理上传的文件，返回解析后的内容和元数据

        Args:
            file_path: 文件路径

        Returns:
            Dict包含：
            - content: 提取的文本内容
            - metadata: 文件元数据
            - chunks: 文本分块（如果适用）
            - file_type: 文件类型
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"文件不存在: {file_path}")

        file_ext = Path(file_path).suffix.lower()

        try:
            if file_ext == '.pdf':
                return self._process_pdf(file_path)
            elif file_ext in ['.docx', '.doc']:
                return self._process_word(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return self._process_excel(file_path)
            elif file_ext in ['.pptx', '.ppt']:
                return self._process_powerpoint(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
                return self._process_image(file_path)
            elif file_ext in ['.txt', '.md', '.json']:
                return self._process_text(file_path)
            else:
                raise ValueError(f"不支持的文件类型: {file_ext}")

        except Exception as e:
            logger.error(f"文件处理失败 {file_path}: {str(e)}")
            raise

    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """处理PDF文件"""
        if not HAS_PYMUPDF:
            raise ImportError("PyMuPDF 未安装，无法处理 PDF 文件")
        
        content = []
        metadata = {}

        try:
            doc = fitz.open(file_path)
            metadata = doc.metadata or {}
            metadata['page_count'] = len(doc)
            metadata['file_size'] = os.path.getsize(file_path)

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    content.append({
                        'page': page_num + 1,
                        'content': text.strip()
                    })

            doc.close()

        except Exception as e:
            logger.error(f"PDF处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join([item['content'] for item in content]),
            'metadata': json.dumps(metadata),
            'chunks': content,
            'file_type': 'pdf'
        }

    def _process_word(self, file_path: str) -> Dict[str, Any]:
        """处理Word文档"""
        if not HAS_DOCX:
            raise ImportError("python-docx 未安装，无法处理 Word 文件")
        
        content = []
        metadata = {}

        try:
            doc = DocxDocument(file_path)
            metadata['paragraph_count'] = len(doc.paragraphs)
            metadata['file_size'] = os.path.getsize(file_path)

            for para in doc.paragraphs:
                if para.text.strip():
                    content.append(para.text.strip())

            # 提取表格内容
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                if table_data:
                    content.append(f"\n表格:\n{json.dumps(table_data, ensure_ascii=False)}")

        except Exception as e:
            logger.error(f"Word处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join(content),
            'metadata': json.dumps(metadata),
            'chunks': [{'content': para, 'type': 'paragraph'} for para in content],
            'file_type': 'word'
        }

    def _process_excel(self, file_path: str) -> Dict[str, Any]:
        """处理Excel文件"""
        if not HAS_OPENPYXL:
            raise ImportError("openpyxl 未安装，无法处理 Excel 文件")
        
        content = []
        metadata = {}

        try:
            workbook = load_workbook(file_path, read_only=True)
            metadata['sheet_count'] = len(workbook.sheetnames)
            metadata['file_size'] = os.path.getsize(file_path)

            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                sheet_content = []
                
                for row in worksheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        sheet_content.append([str(cell) if cell is not None else '' for cell in row])
                
                if sheet_content:
                    content.append(f"工作表: {sheet_name}\n{json.dumps(sheet_content, ensure_ascii=False)}")

            workbook.close()

        except Exception as e:
            logger.error(f"Excel处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join(content),
            'metadata': json.dumps(metadata),
            'chunks': content,
            'file_type': 'excel'
        }

    def _process_powerpoint(self, file_path: str) -> Dict[str, Any]:
        """处理PowerPoint文件"""
        if not HAS_PPTX:
            raise ImportError("python-pptx 未安装，无法处理 PPT 文件")
        
        content = []
        metadata = {}

        try:
            presentation = Presentation(file_path)
            metadata['slide_count'] = len(presentation.slides)
            metadata['file_size'] = os.path.getsize(file_path)

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text.strip())
                
                if slide_content:
                    content.append(f"幻灯片 {slide_num}:\n" + '\n'.join(slide_content))

        except Exception as e:
            logger.error(f"PowerPoint处理失败: {str(e)}")
            raise

        return {
            'content': '\n\n'.join(content),
            'metadata': json.dumps(metadata),
            'chunks': content,
            'file_type': 'powerpoint'
        }

    def _process_text(self, file_path: str) -> Dict[str, Any]:
        """处理文本文件"""
        metadata = {}

        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            metadata['encoding'] = 'utf-8'
        except UnicodeDecodeError:
            try:
                with open(file_path, 'r', encoding='gbk') as f:
                    content = f.read()
                metadata['encoding'] = 'gbk'
            except Exception as e:
                raise ValueError(f"无法解析文本文件编码: {str(e)}")

        metadata['file_size'] = os.path.getsize(file_path)
        paragraphs = [para.strip() for para in content.split('\n\n') if para.strip()]

        return {
            'content': content,
            'metadata': json.dumps(metadata),
            'chunks': [{'content': para, 'type': 'paragraph'} for para in paragraphs],
            'file_type': 'text'
        }

    def _process_image(self, file_path: str) -> Dict[str, Any]:
        """处理图片文件"""
        metadata = {
            'file_size': os.path.getsize(file_path),
            'image_path': file_path,
        }
        
        if HAS_PIL:
            try:
                img = Image.open(file_path)
                metadata['width'] = img.width
                metadata['height'] = img.height
                metadata['format'] = img.format
            except Exception as e:
                logger.warning(f"无法读取图片信息: {e}")

        # 注意：OCR 功能需要额外配置，这里只返回基本信息
        return {
            'content': f"图片文件: {Path(file_path).name}\n文件大小: {metadata['file_size']} 字节",
            'metadata': json.dumps(metadata),
            'chunks': [],
            'file_type': 'image'
        }

    def get_supported_formats(self) -> List[str]:
        """获取支持的文件格式"""
        return ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.md', '.json',
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']

